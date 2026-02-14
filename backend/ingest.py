"""Document ingestion: scan folders, extract text, chunk."""
from pathlib import Path
from dataclasses import dataclass
from typing import Iterator
import logging

from backend.config import settings

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt"}
CHUNK_SIZE = settings.chunk_size
CHUNK_OVERLAP = settings.chunk_overlap


@dataclass
class DocumentChunk:
    """A chunk of document with metadata for RAG."""
    content: str
    source_path: str
    source_name: str
    page_or_sheet: str | None
    last_modified: str
    folder_tag: str
    doc_type: str


def _read_pdf(path: Path) -> list[tuple[str, int]]:
    """Extract text by page from PDF. Returns [(text, page_num), ...]."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        return [(p.extract_text() or "", i + 1) for i, p in enumerate(reader.pages)]
    except Exception as e:
        logger.warning("PDF read failed %s: %s", path, e)
        return []


def _read_docx(path: Path) -> list[tuple[str, int]]:
    """Extract paragraphs. Page numbers approximate by block."""
    try:
        from docx import Document
        doc = Document(path)
        text_by_page: list[str] = []
        current = []
        for p in doc.paragraphs:
            current.append(p.text)
            if len("\n".join(current)) > 1500:  # rough page break
                text_by_page.append("\n".join(current))
                current = []
        if current:
            text_by_page.append("\n".join(current))
        return [(t, i + 1) for i, t in enumerate(text_by_page)] if text_by_page else [("\n".join(p.text for p in doc.paragraphs), 1)]
    except Exception as e:
        logger.warning("DOCX read failed %s: %s", path, e)
        return []


def _read_xlsx(path: Path) -> list[tuple[str, int]]:
    """Extract text from sheet names and cells. Sheet index as 'page'."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        out = []
        for idx, sheet in enumerate(wb.worksheets):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                rows.append(" ".join(str(c) if c is not None else "" for c in row))
            out.append(("\n".join(rows), idx + 1))
        return out if out else [("", 1)]
    except Exception as e:
        logger.warning("XLSX read failed %s: %s", path, e)
        return []


def _read_pptx(path: Path) -> list[tuple[str, int]]:
    """Extract text from slides. Slide number as page."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for i, slide in enumerate(prs.slides):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            out.append(("\n".join(text), i + 1))
        return out if out else [("", 1)]
    except Exception as e:
        logger.warning("PPTX read failed %s: %s", path, e)
        return []


def _read_txt(path: Path) -> list[tuple[str, int]]:
    """Read plain text. Single 'page' or split by double newlines."""
    try:
        raw = path.read_text(encoding="utf-8", errors="replace")
        blocks = [b.strip() for b in raw.split("\n\n") if b.strip()]
        if not blocks:
            return [(raw, 1)]
        return [(b, i + 1) for i, b in enumerate(blocks)]
    except Exception as e:
        logger.warning("TXT read failed %s: %s", path, e)
        return []


def extract_text_by_page(path: Path) -> list[tuple[str, int]]:
    """Dispatch by extension. Returns [(text, page_or_sheet_num), ...]."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".xlsx":
        return _read_xlsx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".txt":
        return _read_txt(path)
    return []


def _chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks by character (can be swapped for sentence-aware)."""
    if not text or not text.strip():
        return []
    text = text.replace("\r", "").strip()
    chunks = []
    start = 0
    while start < len(text):
        end = start + size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
    return chunks


def _folder_tag(path: Path, root: Path) -> str:
    """Folder or department tag from path."""
    try:
        rel = path.relative_to(root)
        parts = rel.parts
        if len(parts) > 1:
            return parts[0]
        return path.parent.name or "Root"
    except ValueError:
        return path.parent.name or "Root"


def iter_documents(root: Path | None = None) -> Iterator[DocumentChunk]:
    """Scan root (default: settings.docs_dir), extract and chunk; yield DocumentChunk."""
    root = root or settings.docs_dir
    if not root.exists():
        logger.warning("Docs root does not exist: %s", root)
        return
    root = root.resolve()
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        try:
            mtime = path.stat().st_mtime
            from datetime import datetime
            last_modified = datetime.fromtimestamp(mtime).strftime("%Y-%m-%d %H:%M")
        except OSError:
            last_modified = ""
        folder_tag = _folder_tag(path, root)
        doc_type = path.suffix.lower().lstrip(".")
        for page_text, page_num in extract_text_by_page(path):
            for chunk in _chunk_text(page_text):
                if not chunk:
                    continue
                yield DocumentChunk(
                    content=chunk,
                    source_path=str(path),
                    source_name=path.name,
                    page_or_sheet=str(page_num),
                    last_modified=last_modified,
                    folder_tag=folder_tag,
                    doc_type=doc_type,
                )
